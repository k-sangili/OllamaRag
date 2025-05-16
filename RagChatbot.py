from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.vectorstores import Pinecone
from langchain.chains import RetrievalQA
from langchain_community.llms import Ollama
from pinecone import Pinecone as PineconeClient
from dotenv import load_dotenv
import os
import gradio as gr
import glob
import pathlib
import PyPDF2
import io
import docx
import pptx
from pptx import Presentation


# Install required packages:
# pip install langchain langchain-community langchain-core sentence-transformers pinecone-client ollama-python python-dotenv gradio PyPDF2 python-docx python-pptx


# Setup Environment Variables (Replace with your info)
load_dotenv()


PINECONE_API_KEY = os.getenv('PINECONE_API_KEY')
PINECONE_ENVIRONMENT = os.getenv('PINECONE_ENVIRONMENT')
PINECONE_INDEX_NAME = os.getenv('PINECONE_INDEX_NAME')


# Validate credentials
if PINECONE_API_KEY == 'PINECONE_API_KEY':
    raise ValueError("Please set your Pinecone API key")
if PINECONE_ENVIRONMENT == 'PINECONE_ENVIRONMENT':
    raise ValueError("Please set your Pinecone environment")
if PINECONE_INDEX_NAME == 'PINECONE_INDEX_NAME':
    raise ValueError("Please set your Pinecone index name")


# Function to load documents from a folder
def load_documents_from_folder(folder_path):
    documents = []
    supported_extensions = ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv', '.pdf', '.docx', '.pptx']
    
    # Get all files with supported extensions
    for ext in supported_extensions:
        files = glob.glob(os.path.join(folder_path, f"*{ext}"))
        for file_path in files:
            try:
                if ext == '.pdf':
                    # Handle PDF files
                    with open(file_path, 'rb') as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        text = ""
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            text += page.extract_text() + "\n"
                        
                        # Add file path as metadata
                        documents.append({
                            "content": text,
                            "metadata": {"source": file_path}
                        })
                
                elif ext == '.docx':
                    # Handle Word documents
                    doc = docx.Document(file_path)
                    text = ""
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                    
                    # Add file path as metadata
                    documents.append({
                        "content": text,
                        "metadata": {"source": file_path}
                    })
                
                elif ext == '.pptx':
                    # Handle PowerPoint presentations
                    presentation = Presentation(file_path)
                    text = ""
                    
                    # Extract text from slides
                    for i, slide in enumerate(presentation.slides):
                        text += f"Slide {i+1}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                        text += "\n"
                    
                    # Add file path as metadata
                    documents.append({
                        "content": text,
                        "metadata": {"source": file_path}
                    })
                
                else:
                    # Handle text-based files
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        # Add file path as metadata
                        documents.append({
                            "content": content,
                            "metadata": {"source": file_path}
                        })
                print(f"Loaded: {file_path}")
            except Exception as e:
                print(f"Error loading {file_path}: {str(e)}")
    
    return documents


# Ingest documents and create embeddings
def process_documents(folder_path):
    # Load documents
    documents = load_documents_from_folder(folder_path)
    
    if not documents:
        return "No documents found in the specified folder."
    
    # Initialize text splitter with smaller chunk size for better metadata management
    text_splitter = CharacterTextSplitter(
        chunk_size=300,  # Smaller chunk size
        chunk_overlap=30,  # 10% overlap
        separator="\n"    # Split preferably at newlines
    )
    
    # Split documents into chunks
    chunks = []
    for doc in documents:
        doc_chunks = text_splitter.split_text(doc["content"])
        # Add source to each chunk as metadata
        chunks.extend([(chunk, doc["metadata"]["source"]) for chunk in doc_chunks])
    
    if not chunks:
        return "No text chunks were created from the documents."
    
    print(f"Created {len(chunks)} chunks from {len(documents)} documents")
    
    # Initialize embeddings
    embeddings = SentenceTransformerEmbeddings(model_name="all-mpnet-base-v2")
    
    try:
        # Initialize Pinecone client
        pc = PineconeClient(api_key=PINECONE_API_KEY, environment=PINECONE_ENVIRONMENT)
        
        # Check if index exists
        if PINECONE_INDEX_NAME not in pc.list_indexes().names():
            return f"Index {PINECONE_INDEX_NAME} does not exist in your Pinecone account"
        
        # Create an index instance
        index = pc.Index(PINECONE_INDEX_NAME)
        
        # Prepare vectors for upsert with metadata - limit metadata size
        vectors = []
        skipped_chunks = 0
        
        for i, (chunk, source) in enumerate(chunks):
            # Create a truncated version of the chunk for metadata to avoid size limit issues
            # Store only the first 500 characters in metadata
            truncated_text = chunk[:500] if len(chunk) > 500 else chunk
            
            # Generate a unique ID for the vector
            vector_id = f"{pathlib.Path(source).stem}-{i}"
            
            try:
                embedding = embeddings.embed_query(chunk)
                
                # Create the vector
                vector = {
                    'id': vector_id,
                    'values': embedding,
                    'metadata': {
                        'text': truncated_text,  # Truncated for metadata
                        'source': os.path.basename(source)  # Just filename, not full path
                    }
                }
                
                vectors.append(vector)
                
                # Upsert in smaller batches to avoid payload size issues
                if len(vectors) >= 50:  # Process in batches of 50
                    index.upsert(vectors=vectors)
                    vectors = []  # Clear after upsert
                    
            except Exception as e:
                print(f"Error processing chunk {i}: {str(e)}")
                skipped_chunks += 1
        
        # Upsert any remaining vectors
        if vectors:
            index.upsert(vectors=vectors)
        
        return f"Successfully processed {len(documents)} documents and created {len(chunks) - skipped_chunks} chunks. Skipped {skipped_chunks} chunks due to errors."
    
    except Exception as e:
        return f"Error during embedding and indexing: {str(e)}"


# Query function
def ask_question(query):
    try:
        # Verify Ollama is running and model is available
        ollama_llm = Ollama(model="deepseek-r1")
        
        # Initialize embeddings
        embeddings = SentenceTransformerEmbeddings(model_name="all-mpnet-base-v2")
        
        # Initialize Pinecone client
        pc = PineconeClient(api_key=PINECONE_API_KEY, environment=PINECONE_ENVIRONMENT)
        index = pc.Index(PINECONE_INDEX_NAME)
        
        # Use Langchain's Pinecone Vectorstore class for retrieval
        vectorstore = Pinecone(index, embeddings.embed_query, 'text')
        
        qa = RetrievalQA.from_chain_type(
            llm=ollama_llm,
            chain_type="stuff",
            retriever=vectorstore.as_retriever(search_kwargs={"k": 5})  # Retrieve top 5 most relevant chunks
        )
        
        # Get response
        response = qa.invoke(query)
        
        # Get retrieved documents for debugging
        retrieved_docs = vectorstore.as_retriever().get_relevant_documents(query)
        
        response_text = f"Response: {response['result']}\n\nRetrieved Documents:\n"
        for i, doc in enumerate(retrieved_docs):
            source = doc.metadata.get('source', 'Unknown source')
            text_preview = doc.page_content if len(doc.page_content) <= 100 else doc.page_content[:100] + "..."
            response_text += f"{i+1}. Source: {source}\n   Preview: {text_preview}\n\n"
        
        return response_text
    
    except Exception as e:
        return f"Error during retrieval and generation: {str(e)}"


# Create Gradio Interface with two tabs
with gr.Blocks(title="RAG Chatbot with Local Documents") as app:
    gr.Markdown("# RAG Chatbot with Ollama and Pinecone")
    
    with gr.Tabs():
        with gr.Tab("Upload Documents"):
            folder_input = gr.Textbox(
                label="Enter folder path containing documents",
                placeholder="/path/to/your/documents"
            )
            file_types = gr.Markdown("""
            Supported file types:
            - Text files (.txt)
            - Markdown files (.md)
            - Code files (.py, .js, .html, .css)
            - Data files (.json, .csv)
            - PDF files (.pdf)
            """)
            process_btn = gr.Button("Process Documents")
            process_output = gr.Textbox(label="Processing Results")
            
            process_btn.click(
                fn=process_documents,
                inputs=folder_input,
                outputs=process_output
            )
        
        with gr.Tab("Ask Questions"):
            question_input = gr.Textbox(
                lines=2,
                label="Enter your question:",
                placeholder="What would you like to know about your documents?"
            )
            answer_btn = gr.Button("Ask")
            answer_output = gr.Textbox(lines=10, label="Response")
            
            answer_btn.click(
                fn=ask_question,
                inputs=question_input,
                outputs=answer_output
            )


# Launch the Interface
if __name__ == "__main__":
    app.launch()